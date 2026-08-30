"""
Assembles design/ui-kit/*.png into design/emqnote-ui-kit.pptx — a parts bin, not a mockup.

    python3 -m venv .venv && .venv/bin/pip install python-pptx
    npm run ui:deck                       # or: .venv/bin/python scripts/build-ui-deck.py
    npm run ui:deck -- --preview /tmp/deck # also write an HTML rendering to look at

Run `npm run ui:kit` first: this reads design/ui-kit/manifest.json and nothing else about
the app, except the colour tokens, which it parses straight out of src/renderer/styles.css
so the palette slide cannot drift from the stylesheet.

Python rather than TypeScript, and python-pptx rather than hand-written OOXML, for one
reason: a deck PowerPoint refuses to open is a worse failure than a dependency in a venv,
and there is no PowerPoint here to find that out with. Nothing is added to package.json —
`dependencies` there stays minimal and `check:bundle` never sees this.

What it checks before it exits, since the deck itself cannot be opened here: every shape
inside the slide, no two pictures overlapping, every XML part well formed, every media part
referenced. A non-zero exit means one of those failed.
"""
import json
import os
import re
import sys
import xml.dom.minidom
import zipfile
from pathlib import Path

# `npm run ui:deck` calls plain `python3`, which is the interpreter that will not have
# python-pptx: it lives in the venv the docstring asks for. Rather than fail with the right
# advice and the wrong interpreter, hand over to the venv's python once — the guard keeps
# that from becoming a loop when the venv is the one that is missing the package.
if "PPTX_HANDED_OVER" not in os.environ:
    venv = Path(".venv/bin/python")
    if venv.exists():
        try:
            import pptx  # noqa: F401
        except ModuleNotFoundError:
            os.execve(str(venv), [str(venv), *sys.argv], {**os.environ, "PPTX_HANDED_OVER": "1"})

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt
except ModuleNotFoundError:
    sys.exit(
        "python-pptx is not installed. It is deliberately not a project dependency:\n"
        "  python3 -m venv .venv && .venv/bin/pip install python-pptx\n"
        "  .venv/bin/python scripts/build-ui-deck.py"
    )

ROOT = Path(".")
KIT = ROOT / "design/ui-kit"
OUT = ROOT / "design/emqnote-ui-kit.pptx"

SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN_X, TOP, BOTTOM = 0.42, 1.12, 7.14
GAP_X, GAP_Y, CAPTION_H = 0.22, 0.34, 0.30

# One scale for the whole kit, and the reason is assembly rather than fit: parts only
# compose with each other if they are all drawn at the same size, and a mockup wants the
# window shell and the rows that go on it to agree. Half of the app's pixels puts the
# 1440-wide library window on a 13.3-inch slide with room to spare, and the parts are
# captured at 3× — so one on a slide is still six times oversampled and stays sharp
# however far it is zoomed or printed.
KIT_SCALE = 0.5

PAGE_BG = RGBColor.from_string("EEF0F3")
INK = RGBColor.from_string("1B1C1F")
MUTED = RGBColor.from_string("6B7079")
FONT = "Segoe UI"

FAMILY_ORDER = [
    ("Windows", "The two windows, whole"),
    ("Shells", "Blanks to build on — empty, ready for parts"),
    ("Chrome", "Title bars and footers"),
    ("Folder tree", "The left pane: where"),
    ("Note list", "The middle pane: which"),
    ("Reader", "The right pane: what"),
    ("Header block", "The note's own fields, in both windows"),
    ("Editor", "What the editor draws — the type specimens"),
    ("Menus", "Menus"),
    ("Dialogs", "Dialogs and panels"),
    ("Views", "Tasks, filters, files"),
]


def themes():
    """The two themes' roles, read out of styles.css so this cannot drift from it."""
    css = (ROOT / "src/renderer/styles.css").read_text(encoding="utf8")
    light = re.search(r"@media \(prefers-color-scheme: light\).*?:root \{(.*?)\}", css, re.S)
    dark = re.search(r"^:root \{(.*?)\}", css, re.S | re.M)
    pairs = lambda text: dict(re.findall(r"(--[a-z-]+):\s*([^;]+);", text))
    return pairs(light.group(1)), pairs(dark.group(1))


def add_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = PAGE_BG
    background.line.fill.background()
    background.shadow.inherit = False
    background.name = "slide background"

    box = slide.shapes.add_textbox(
        Inches(MARGIN_X), Inches(0.34), Inches(SLIDE_W - 2 * MARGIN_X), Inches(0.5)
    )
    frame = box.text_frame
    frame.word_wrap = True
    run = frame.paragraphs[0].add_run()
    run.text = title
    run.font.size, run.font.bold, run.font.name = Pt(22), True, FONT
    run.font.color.rgb = INK
    if subtitle is not None:
        run = frame.add_paragraph().add_run()
        run.text = subtitle
        run.font.size, run.font.name = Pt(11), FONT
        run.font.color.rgb = MUTED
    return slide


def caption(slide, left, top, width, text):
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Inches(CAPTION_H))
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.size, run.font.name = Pt(8.5), FONT
    run.font.color.rgb = MUTED
    box.name = "caption"
    return box


def place(prs, family, blurb, items):
    """Shelf-packs a family's parts across as many slides as it takes.

    Tallest first, and stable, so equal-sized neighbours keep the order they were captured
    in — the four header cells stay When, Tags, Where, Who. Without the sort a pane four
    inches tall sets the shelf height and everything short after it leaves a band of empty
    slide; with it, the short parts fill in beside the tall one.
    """
    items = sorted(items, key=lambda item: -item["cssHeight"])
    slide = add_slide(prs, family, blurb)
    notes = []
    x, y, shelf = Inches(MARGIN_X), Inches(TOP), 0
    limit_x, limit_y = Inches(SLIDE_W - MARGIN_X), Inches(BOTTOM)
    max_w, max_h = Inches(SLIDE_W - 2 * MARGIN_X), Inches(BOTTOM - TOP - CAPTION_H)

    for item in items:
        native_w = Inches(item["cssWidth"] / 96 * KIT_SCALE)
        native_h = Inches(item["cssHeight"] / 96 * KIT_SCALE)
        scale = min(1.0, max_w / native_w, max_h / native_h)
        width, height = int(native_w * scale), int(native_h * scale)

        if x + width > limit_x and x > Inches(MARGIN_X):                # next shelf
            x, y, shelf = Inches(MARGIN_X), y + shelf + Inches(CAPTION_H + GAP_Y), 0
        # `y > TOP`, not `shelf > 0`: after a wrap the shelf is empty again, and guarding on
        # it put a tall part below the bottom of the slide instead of on the next one.
        if y + height + Inches(CAPTION_H) > limit_y and y > Inches(TOP):  # next slide
            slide.notes_slide.notes_text_frame.text = "\n".join(notes)
            notes = []
            slide = add_slide(prs, family, f"{blurb} (continued)")
            x, y, shelf = Inches(MARGIN_X), Inches(TOP), 0

        picture = slide.shapes.add_picture(str(KIT / item["file"]), x, y, width=width, height=height)
        picture.name = item["selector"]
        label = f"{item['caption']} · {item['cssWidth']}×{item['cssHeight']} px"
        if scale < 0.999:
            label += f" · smaller than kit scale ({round(scale * KIT_SCALE * 100)}%)"
        caption(slide, x, y + height + Inches(0.05), min(max(width, Inches(1.6)), limit_x - x), label)
        notes.append(
            f"{item['file']} — {item['selector']} — {item['cssWidth']}×{item['cssHeight']} px at 3×"
        )

        x += width + Inches(GAP_X)
        shelf = max(shelf, height)

    slide.notes_slide.notes_text_frame.text = "\n".join(notes)


def cover(prs, count):
    slide = add_slide(
        prs, "emqnote — UI kit",
        "Every part of the app, photographed from the running app. Light theme.",
    )
    lines = [
        ("How to use it", True),
        ("Copy a part off these slides onto your own, and arrange. Text goes on top in a text box —"
         " a part is a picture and cannot be re-typed inside.", False),
        ("Two blanks are on the “Shells” slide: an empty library window and an empty capture window,"
         " full size. Build on those.", False),
        ("Sizes", True),
        ("Every part is at the same scale — half of the app's own pixels — so parts fit together"
         " without any adjusting. The captions give the true pixel size. Hold Shift when you resize,"
         " and resize a whole group rather than one part, or the parts stop matching.", False),
        ("Captured at 3×, so a part stays sharp however far you zoom it or print it.", False),
        ("Where each part comes from", True),
        ("A part's CSS selector is its shape name — Home ▸ Select ▸ Selection Pane shows it. The same"
         " list, with pixel sizes, is in each slide's notes. Hand those to a coding agent and it knows"
         " exactly which element you moved.", False),
        ("A snapshot, not a source", True),
        ("Captured from the app as it stood when this was built. Change the renderer and this deck is"
         " out of date — it does not update itself. `npm run ui:kit` then `npm run ui:deck` rebuilds"
         " it.", False),
    ]
    box = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(1.35), Inches(8.4), Inches(5.4))
    frame = box.text_frame
    frame.word_wrap = True
    first = True
    for text, heading in lines:
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(13 if heading else 11)
        run.font.bold = heading
        run.font.color.rgb = INK if heading else MUTED
        paragraph.space_before = Pt(12 if heading else 3)

    slide.shapes.add_picture(
        str(KIT / "window-library.png"), Inches(8.95), Inches(1.5), width=Inches(3.95)
    ).name = "cover picture"
    slide.notes_slide.notes_text_frame.text = (
        f"{count} parts, captured from the running app under Xvfb by scripts/export-ui-kit.ts."
    )


def palette(prs, light, dark):
    slide = add_slide(
        prs, "Colour roles",
        "Six roles carry the whole UI (B87). These are real shapes — copy one, or lift its fill.",
    )
    roles = [
        ("--background", "the page"), ("--surface", "chrome and panels"),
        ("--field", "a box you type in"), ("--border", "any divider"),
        ("--text", "words"), ("--muted", "quieter words"),
        ("--accent", "links and focus"), ("--warning", "something is off"),
        ("--highlight", "==highlight=="),
    ]
    for row, (name, tokens) in enumerate(
        (("Light — what this kit is drawn in", light), ("Dark", dark))
    ):
        top = Inches(1.3 + row * 2.75)
        label = slide.shapes.add_textbox(Inches(MARGIN_X), top, Inches(6), Inches(0.3))
        run = label.text_frame.paragraphs[0].add_run()
        run.text = name
        run.font.name, run.font.size, run.font.bold = FONT, Pt(12), True
        run.font.color.rgb = INK
        label.name = "palette heading"

        for column, (role, what) in enumerate(roles):
            value = tokens.get(role, "").strip()
            if not value.startswith("#"):
                continue
            left = Inches(MARGIN_X + column * 1.42)
            chip = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top + Inches(0.38), Inches(1.24), Inches(1.0)
            )
            chip.fill.solid()
            chip.fill.fore_color.rgb = RGBColor.from_string(value.lstrip("#").upper()[:6])
            chip.line.color.rgb = RGBColor.from_string("C7CCD4")
            chip.line.width = Pt(0.75)
            chip.shadow.inherit = False
            chip.name = role
            chip.text_frame.text = ""
            caption(slide, left, top + Inches(1.44), Inches(1.3), f"{role}\n{value}\n{what}")

    slide.notes_slide.notes_text_frame.text = (
        "Parsed from src/renderer/styles.css at build time — the :root block and the "
        "prefers-color-scheme: light block. Six roles, declared once per theme (B87)."
    )


def check(preview_dir):
    """Reads the finished deck back, and optionally draws it as HTML so it can be looked at."""
    prs = Presentation(str(OUT))
    problems, pictures, drawn_slides = [], 0, []

    for number, slide in enumerate(prs.slides, start=1):
        boxes, drawn = [], []
        for shape in slide.shapes:
            left, top = shape.left or 0, shape.top or 0
            width, height = shape.width or 0, shape.height or 0
            if shape.name != "slide background" and (
                left < 0 or top < 0 or left + width > prs.slide_width or top + height > prs.slide_height
            ):
                problems.append(f"slide {number}: {shape.name} runs off the slide")
            if shape.shape_type == 13:  # PICTURE
                pictures += 1
                for other, box in boxes:
                    if not (left + width <= box[0] or box[0] + box[2] <= left
                            or top + height <= box[1] or box[1] + box[3] <= top):
                        problems.append(f"slide {number}: {shape.name} overlaps {other}")
                boxes.append((shape.name, (left, top, width, height)))
                drawn.append(("img", left, top, width, height, shape.image.blob, shape.name))
            else:
                fill = None
                try:
                    fill = str(shape.fill.fore_color.rgb)
                except Exception:
                    pass
                text = shape.text_frame.text if shape.has_text_frame else ""
                kind = "bg" if shape.name == "slide background" else "shape"
                drawn.append((kind, left, top, width, height, fill, text))
        drawn_slides.append(drawn)

    with zipfile.ZipFile(OUT) as deck:
        names = deck.namelist()
        for part in [n for n in names if n.endswith((".xml", ".rels"))]:
            try:
                xml.dom.minidom.parseString(deck.read(part))
            except Exception as error:
                problems.append(f"{part} is not well formed: {error}")
        references = b"".join(deck.read(n) for n in names if n.endswith(".rels"))
        for media in [n for n in names if n.startswith("ppt/media/")]:
            if media.split("/")[-1].encode() not in references:
                problems.append(f"{media} is in the package but nothing points at it")

    if preview_dir is not None:
        stage = Path(preview_dir)
        stage.mkdir(parents=True, exist_ok=True)
        scale = 900 / prs.slide_width
        html = ['<style>body{background:#33363b;margin:0;font-family:"Segoe UI",sans-serif}'
                '.slide{position:relative;margin:14px auto;background:#fff;overflow:hidden}</style>']
        for drawn in drawn_slides:
            html.append(f'<div class="slide" style="width:{prs.slide_width * scale:.0f}px;'
                        f'height:{prs.slide_height * scale:.0f}px">')
            for item in drawn:
                kind = item[0]
                style = (f"position:absolute;left:{item[1] * scale:.1f}px;top:{item[2] * scale:.1f}px;"
                         f"width:{item[3] * scale:.1f}px;height:{item[4] * scale:.1f}px;")
                if kind == "img":
                    name = f"{abs(hash(item[5])):x}.png"
                    (stage / name).write_bytes(item[5])
                    html.append(f'<img src="{name}" style="{style}">')
                elif kind == "bg":
                    html.append(f'<div style="{style}background:#EEF0F3"></div>')
                else:
                    fill = f"background:#{item[5]};" if item[5] else ""
                    body = (item[6] or "").replace("&", "&amp;").replace("<", "&lt;").replace("\n", "<br>")
                    html.append(f'<div style="{style}{fill}font-size:9px;color:#333;'
                                f'line-height:1.25;overflow:hidden">{body}</div>')
            html.append("</div>")
        (stage / "preview.html").write_text("\n".join(html), encoding="utf8")
        print(f"preview: {stage / 'preview.html'}")

    print(f"{len(prs.slides._sldIdLst)} slides, {pictures} pictures")
    if problems:
        print("\n".join(problems))
    else:
        print("every shape is inside its slide, no two pictures overlap, "
              "every XML part is well formed and every picture is referenced")
    return 1 if problems else 0


def main():
    if not (KIT / "manifest.json").exists():
        sys.exit("design/ui-kit/manifest.json is missing — run `npm run ui:kit` first")
    manifest = json.loads((KIT / "manifest.json").read_text(encoding="utf8"))
    light, dark = themes()

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SLIDE_W), Inches(SLIDE_H)

    cover(prs, len(manifest))
    palette(prs, light, dark)

    filed = set()
    for family, blurb in FAMILY_ORDER:
        items = [item for item in manifest if item["family"] == family]
        filed.update(item["file"] for item in items)
        if items:
            place(prs, family, blurb, items)
    left_over = [item for item in manifest if item["file"] not in filed]
    if left_over:
        place(prs, "Everything else", "Not filed under a family", left_over)

    prs.save(str(OUT))
    print(f"{OUT}: {len(manifest)} parts, {OUT.stat().st_size / 1048576:.1f} MB")

    preview = None
    if "--preview" in sys.argv:
        preview = sys.argv[sys.argv.index("--preview") + 1]
    return check(preview)


sys.exit(main())
